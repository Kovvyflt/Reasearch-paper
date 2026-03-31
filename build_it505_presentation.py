#!/usr/bin/env python3
"""Generate IT505 NorthStar Digital project presentation from report sources."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT_PATH = "/workspace/IT505_NorthStar_Digital_Project_Presentation.pptx"

ACCENT = RGBColor(0x1A, 0x36, 0x5D)  # deep blue
SUB = RGBColor(0x44, 0x44, 0x44)


def set_slide_title(slide, text):
    if slide.shapes.title:
        slide.shapes.title.text = text
        p = slide.shapes.title.text_frame.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = ACCENT


def add_bullets(slide, lines, left=0.5, top=1.35, width=9.0, height=5.5):
    """Add or replace body with bullet list."""
    body = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            body = shape
            break
    if body is None:
        body = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(14 if len(lines) > 8 else 15)
        p.font.color.rgb = SUB
        p.space_after = Pt(6)


def add_two_column_bullets(slide, left_lines, right_lines, title_left="In scope", title_right="Out of scope"):
    """Title + two text boxes for comparison (caller sets slide title)."""
    # Remove default placeholder body if present
    for shape in list(slide.shapes):
        if shape.has_text_frame and shape != slide.shapes.title:
            sp = shape.element
            sp.getparent().remove(sp)
    # Left column
    box1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.25), Inches(4.5), Inches(5.7))
    tf1 = box1.text_frame
    p = tf1.paragraphs[0]
    p.text = title_left
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = ACCENT
    for line in left_lines:
        p = tf1.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(13)
        p.font.color.rgb = SUB
        p.space_after = Pt(4)
    # Right column
    box2 = slide.shapes.add_textbox(Inches(5.1), Inches(1.25), Inches(4.3), Inches(5.7))
    tf2 = box2.text_frame
    p = tf2.paragraphs[0]
    p.text = title_right
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = ACCENT
    for line in right_lines:
        p = tf2.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(13)
        p.font.color.rgb = SUB
        p.space_after = Pt(4)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Title ---
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(8.6), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "NorthStar Digital"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "Personal Health Companion App (MVP)"
    p.font.size = Pt(26)
    p.font.color.rgb = SUB
    p.alignment = PP_ALIGN.CENTER
    sub = slide.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(8.6), Inches(2.2))
    stf = sub.text_frame
    stf.paragraphs[0].text = "IT505 Group Project — Project Management Plan"
    stf.paragraphs[0].font.size = Pt(18)
    stf.paragraphs[0].alignment = PP_ALIGN.CENTER
    stf.paragraphs[0].font.color.rgb = SUB
    p = stf.add_paragraph()
    p.text = "Pooja Sharma · Oji Nkechi · Oluwakayode Soyinka · Dominick Tressler"
    p.font.size = Pt(15)
    p.alignment = PP_ALIGN.CENTER
    p = stf.add_paragraph()
    p.text = "Concordia University of Edmonton · March 31, 2026"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER

    # --- Agenda ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Presentation overview")
    add_bullets(
        slide,
        [
            "Background, objectives, and constraints",
            "MVP scope, WBS, and change control",
            "Schedule (Apr 1 – Jul 31, 2026) and Agile delivery",
            "Risk analysis — six prioritized risks with responses",
            "Stakeholder register, power/interest matrix, and communications",
            "Key conclusions",
        ],
    )

    # --- Background ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Background & strategic context")
    add_bullets(
        slide,
        [
            "NorthStar Digital Solutions — growing startup; mobile productivity and wellness apps.",
            "Next growth area: digital health — users increasingly manage health on smartphones.",
            "MVP goal: launch quickly, test demand, learn from early adopters (especially older adults).",
            "Product: simple, intuitive app for health tracking, medications, and emergency contacts.",
        ],
    )

    # --- Objectives & constraints ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Objectives & project constraints")
    add_bullets(
        slide,
        [
            "Deliver a mobile health companion within four months (Apr 1 – Jul 31, 2026).",
            "Give users simple recording of health data and emergency contacts.",
            "Ship an MVP to the market and gather structured feedback.",
            "Demonstrate NorthStar’s ability to deliver a reliable health-related mobile solution.",
            "Constraints: $200,000 budget · 4 months · small cross-functional team · Agile.",
        ],
    )

    # --- MVP features ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "MVP — Core features (initial release)")
    add_bullets(
        slide,
        [
            "Health tracker — blood pressure, daily exercise, weight, well-being notes.",
            "Medication reminder — schedule, push notifications, medication notes.",
            "Medical contacts directory — doctors, clinics, pharmacies, emergency contacts.",
            "Emergency information card — allergies, conditions, blood type, emergency contacts; quick access in emergencies.",
        ],
    )

    # --- Design requirements ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Design requirements")
    add_bullets(
        slide,
        [
            "Simple to use; accessible to older adults; secure and privacy-conscious.",
            "Android and iOS compatibility.",
            "Clear navigation, large readable fonts, minimal setup steps.",
            "Scope emphasizes usability and trust — not feature overload.",
        ],
    )

    # --- Budget ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Budget allocation ($200,000 total)")
    add_bullets(
        slide,
        [
            "Project management — $25,000",
            "Mobile development — $90,000",
            "UI/UX design — $25,000",
            "Quality assurance — $25,000",
            "Marketing & launch — $20,000",
            "Contingency reserve — $15,000",
        ],
    )

    # --- Success criteria ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Project success criteria")
    add_bullets(
        slide,
        [
            "Complete within 4 months and within $200,000.",
            "Core features work on both Android and iOS.",
            "Meet usability and accessibility expectations.",
            "No critical security or privacy gaps at release.",
            "MVP launched and user feedback collected.",
        ],
    )

    # --- Approach ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Project approach (Agile)")
    add_bullets(
        slide,
        [
            "Iterative development with continuous feedback across the lifecycle.",
            "Phases: initiation & planning → requirements & scope → UI/UX design/prototyping → incremental core development → QA → MVP launch & feedback.",
            "Small cross-functional team; regular reviews and sprint cycles.",
            "Project Manager: Oji Nkechi — charter dated March 26, 2026.",
        ],
    )

    # --- Scope ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_title(slide, "Scope statement (summary)")
    add_two_column_bullets(
        slide,
        [
            "Reliable, simple health app for Android/iOS in four months.",
            "Health tracker (BP, exercise, weight, notes).",
            "Medication reminder (schedule, notifications, notes).",
            "Medical contacts & emergency card (allergies, blood type, conditions).",
            "Large fonts, clear navigation, secure privacy-conscious architecture.",
        ],
        [
            "Advanced diagnostic tools.",
            "Wearable hardware integration (e.g. smartwatches).",
            "Desktop and web versions.",
        ],
    )

    # --- WBS ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Work breakdown structure (WBS)")
    add_bullets(
        slide,
        [
            "1 — Project management: initiation (charter, stakeholders); sprint planning & backlog.",
            "2 — UI/UX: interface design (accessibility, large type); prototypes & navigation flow.",
            "3 — App development: health tracker; medication reminder & notifications; contacts & emergency card DB.",
            "4 — QA: cross-platform testing; accessibility & usability testing.",
            "5 — Deployment & launch: App Store/Google Play submission; early-adopter feedback.",
        ],
    )

    # --- Verification & change ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Deliverable verification & scope change")
    add_bullets(
        slide,
        [
            "Verification: phase/sprint reviews; compare to acceptance criteria; functional, usability, compatibility testing; stakeholder sign-off; document acceptance or revisions.",
            "Changes: document new requests in the product backlog; PM assesses impact on $200k budget and 4-month timeline.",
            "Approved changes may defer lower-priority MVP items to a later version to protect the release date.",
        ],
    )

    # --- Schedule ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "High-level schedule (4 months)")
    add_bullets(
        slide,
        [
            "Month 1 (April 2026): Project initiation; requirements; backlog; UI/UX foundations & early prototypes.",
            "Month 2 (May): Complete UI/UX; sprint development — health tracker & medication modules.",
            "Month 3 (June): Contacts & emergency card; integration; begin structured QA cycles.",
            "Month 4 (July): Full QA & accessibility; App Store submissions ≥2 weeks before deadline; MVP launch & feedback collection.",
            "Note: Submit to stores early to allow Apple/Google review cycles; use contingency reserve if needed.",
        ],
    )

    # --- Risk methodology ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Risk management approach")
    add_bullets(
        slide,
        [
            "PMI/PMBOK: individual project risks (resources, time) vs. overall project risk (environment).",
            "Qualitative analysis (probability × impact, 1–5 scales) — appropriate for MVP planning stage.",
            "Score bands: 15–25 Critical · 10–14 High · 5–9 Moderate · 1–4 Low.",
            "Keep the risk register current via sprint retrospectives and phase-gate reviews.",
        ],
    )

    # --- Risk summary table ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Risk register — priority summary")
    add_bullets(
        slide,
        [
            "R-01 Data security breach — Technical — P×I → 15 — CRITICAL — Mitigate",
            "R-02 Regulatory non-compliance (PIPEDA/HIPAA) — Legal — 16 — HIGH — Mitigate",
            "R-03 Personnel skill gap / understaffing — Resource — 16 — HIGH — Mitigate/Accept",
            "R-04 Compressed timeline / schedule overrun — Schedule — 12 — HIGH — Mitigate",
            "R-05 Low-quality MVP / excessive defects — Quality — 9 — MODERATE — Mitigate",
            "R-06 Limited market validation (seniors vs. beta testers) — Market — 12 — HIGH — Mitigate",
        ],
    )

    # --- Critical & legal risks ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Top risks — security & compliance")
    add_bullets(
        slide,
        [
            "R-01 Data security breach: PHI on device/server — risk of unauthorized access or interception.",
            "Response: E2E encryption; penetration testing; RBAC; incident response plan before MVP launch. Owner: Mobile Dev / QA Lead.",
            "R-02 Regulatory: PHI collection — PIPEDA/HIPAA exposure; fines, legal action, or app cessation.",
            "Response: legal/privacy consultant early; privacy impact assessment (PIA); consent & data deletion in MVP. Owner: Project Manager.",
        ],
    )

    # --- Resource, schedule, quality ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Risks — people, schedule, quality, market")
    add_bullets(
        slide,
        [
            "R-03: Hire early; cross-train; agile reviews; contingency contractors (up to $15k reserve). PM owns.",
            "R-04: Timeboxed sprints; de-prioritize non-core work; submit stores 2+ weeks early. PM owns.",
            "R-05: Definition of Done with defect thresholds; QA sprints; Appium-style automation; QA in sprint reviews. QA Lead owns.",
            "R-06: Partner with seniors’ organizations; WCAG 2.1 AA; marketing recruits representative beta testers. Marketing/UI-UX.",
        ],
    )

    # --- Risk conclusion ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Risk analysis — conclusions")
    add_bullets(
        slide,
        [
            "Security and compliance carry the highest long-term consequence (legal, financial, reputation).",
            "Internal delivery risks (personnel, schedule) are acute in a small Agile team with limited redundancy.",
            "Market validation affects roadmap after MVP more than day-one delivery — still plan for representative older-adult testing.",
        ],
    )

    # --- Stakeholders ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Stakeholder register (10 stakeholders)")
    add_bullets(
        slide,
        [
            "NorthStar leadership — sponsor; on-time/budget launch; bi-weekly executive reports.",
            "PM (Oji Nkechi) — delivery; sprint coordination; central communication hub.",
            "Mobile developer — stable requirements; sprint planning from day one; clear acceptance criteria.",
            "UI/UX designer — older-adult usability; early usability check-ins; personas & wireframes.",
            "QA lead — time to test; sprint reviews; cross-platform test plans.",
            "Marketing rep — audience fit; launch & MVP feedback plan.",
            "End users (incl. older adults) — simplicity, privacy; beta & in-app feedback.",
            "Healthcare professionals — clarity it is a personal tool, not a medical device; optional design review.",
            "App stores (Apple/Google) — compliance with privacy/security rules; early guideline review.",
            "Regulators (privacy) — PIPEDA-aligned handling; PIA; minimize data; document decisions.",
        ],
    )

    # --- Matrix ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Power / interest matrix")
    add_bullets(
        slide,
        [
            "Manage closely — Project Manager, end users (high interest; users indirect but essential).",
            "Keep satisfied — App stores, regulatory bodies, NorthStar leadership (high influence; compliance & funding).",
            "Keep informed — Mobile developer, UI/UX, QA, marketing (core team).",
            "Monitor — Healthcare professionals (lower influence; endorsement still valuable).",
        ],
    )

    # --- Communications ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Stakeholder communication plan")
    add_bullets(
        slide,
        [
            "Leadership: executive status report + email — bi-weekly — Project Manager.",
            "Project team: sprint planning, stand-ups, retrospectives — daily / per sprint — PM.",
            "End users: beta program, in-app feedback — post-MVP launch — Marketing representative.",
            "App stores: guideline compliance embedded in schedule; submission milestones tracked early.",
            "Principle: right information, right audience, right time — avoid unnecessary noise.",
        ],
    )

    # --- Conclusion ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_title(slide, "Key takeaways")
    add_bullets(
        slide,
        [
            "Clear MVP scope, WBS, and Agile cadence align delivery with a fixed 4-month budget window.",
            "Privacy, security, and regulation are prioritized alongside schedule and quality risks.",
            "Stakeholder and communication plans keep sponsors, stores, regulators, and users appropriately engaged.",
            "Next steps: execute sprints, maintain the risk register, and submit to app stores with buffer before Jul 31, 2026.",
        ],
    )

    # --- Closing ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank you"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "Questions?"
    p.font.size = Pt(22)
    p.font.color.rgb = SUB
    p.alignment = PP_ALIGN.CENTER

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
